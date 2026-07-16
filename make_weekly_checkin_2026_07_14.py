from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
OUT = ROOT / "7_14_26_Weekly_Check_in.pptx"
prs = Presentation(ROOT / "weekly_checkin_2026_07_07" / "7_7_26_Weekly_Check_in.pptx")

NAVY = RGBColor(22, 28, 36)
BLUE = RGBColor(25, 100, 185)
TEAL = RGBColor(12, 132, 112)
ORANGE = RGBColor(232, 119, 34)
GREY = RGBColor(92, 102, 112)
LINE = RGBColor(220, 226, 232)
PALE = RGBColor(247, 250, 252)


def clear(slide, keep=2):
    for shape in list(slide.shapes)[keep:]:
        slide.shapes._spTree.remove(shape._element)


def add_text(slide, x, y, w, h, value, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_box(slide, x, y, w, h, fill=PALE, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape


def add_picture(slide, filename, x, y, w, h):
    return slide.shapes.add_picture(str(ROOT / filename), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def heading(slide, kicker, title, number):
    add_text(slide, .65, .35, 6, .25, kicker, 8, BLUE, True)
    add_text(slide, .65, .62, 9, .55, title, 28, NAVY, True)
    add_box(slide, .65, 1.42, 12, .01, LINE, LINE)
    add_text(slide, 12.35, 7.02, .35, .18, f"{number}/6", 8, GREY)


def add_tag(slide, x, y, value, color):
    tag = add_box(slide, x, y, 1.25, .32, color, color)
    tag.text_frame.text = value
    paragraph = tag.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.name = "Arial"
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)


# Slide 1
slide = prs.slides[0]
clear(slide, 0)
add_text(slide, .65, .58, 7.5, .55, "7/14/26 Weekly Check-in", 34, NAVY, True)
add_text(slide, .68, 1.28, 6, .35, "Tank scan planner progress", 15, GREY)
add_text(slide, .72, 2.18, 5.1, 1.8, "Before → after mission previews\nCircular sweeps are now closer together\nHorizontal rows are the next interior update", 17)
add_text(slide, .72, 5.55, 1.9, .35, "3", 23, BLUE, True)
add_text(slide, .72, 5.97, 1.9, .24, "tank sizes compared", 8.5, GREY)
add_text(slide, 2.55, 5.55, 1.9, .35, "2", 23, TEAL, True)
add_text(slide, 2.55, 5.97, 2.2, .24, "circular rows maximum", 8.5, GREY)
add_box(slide, 6.2, .62, 6.5, 5.92, RGBColor(255, 255, 255), LINE)
add_picture(slide, "dxf_prediction_validation/checkin/completed_150ft.png", 6.49, .62, 5.92, 5.92)
add_text(slide, 6.25, 6.67, 5, .25, "150ft DXF-completion preview", 10, GREY)
add_text(slide, 12.35, 7.02, .35, .18, "1/6", 8, GREY)

# Slide 2: DXF completion panels
slide = prs.slides[1]
clear(slide, 0)
heading(slide, "DXF COMPLETION OUTPUTS", "Completed layouts for all tank sizes", 2)
after = [
    ("24ft", "dxf_prediction_validation/checkin/completed_24ft.png", BLUE, "completed DXF preview"),
    ("65ft", "dxf_prediction_validation/checkin/completed_65ft.png", TEAL, "completed DXF preview"),
    ("150ft", "dxf_prediction_validation/checkin/completed_150ft.png", ORANGE, "completed DXF preview"),
]
for index, (tank, filename, color, note) in enumerate(after):
    x = .55 + index * 4.18
    add_box(slide, x, 1.72, 3.75, 3.75, RGBColor(255, 255, 255), LINE)
    add_picture(slide, filename, x, 1.72, 3.75, 3.75)
    add_tag(slide, x + .22, 1.95, tank, color)
    add_text(slide, x, 5.72, 3.7, .28, note, 11, GREY, True, PP_ALIGN.CENTER)
add_text(slide, .75, 6.35, 11.2, .35, "Each completed layout combines the observed circular-sweep weld fragments with predicted missing plate and weld geometry.", 15, NAVY, True)
add_text(slide, 12.35, 7.02, .35, .18, "2/6", 8, GREY)

# Slide 3: before/after summary
slide = prs.slides[2]
clear(slide, 0)
heading(slide, "BEFORE → AFTER", "What changed in the planner", 3)
for x, label, width in [(.75, "tank", 1.5), (2.35, "before", 3.0), (6.0, "after / current", 3.7), (10.1, "status", 2.0)]:
    add_box(slide, x, 1.75, width, .42, RGBColor(230, 239, 246), RGBColor(230, 239, 246))
    add_text(slide, x + .12, 1.84, width - .2, .22, label, 11, BLUE, True)
rows = [
    ("24ft", "Best mission plan preview", "Saved current reference", "Horizontal rows not present"),
    ("65ft", "Best mission plan preview", "Weld-bounded current preview", "Horizontal rows not present"),
    ("150ft", "Best mission plan preview", "On-segment preview", "Horizontal rows not yet implemented"),
]
for index, row in enumerate(rows):
    y = 2.3 + index * .85
    if index % 2 == 0:
        add_box(slide, .75, y, 11.35, .68, PALE, PALE)
    add_text(slide, .87, y + .18, .9, .25, row[0], 13, NAVY, True)
    add_text(slide, 2.47, y + .12, 3.25, .4, row[1], 11, GREY)
    add_text(slide, 6.12, y + .12, 3.8, .4, row[2], 11, GREY)
    add_text(slide, 10.22, y + .12, 1.75, .4, row[3], 10, ORANGE, True)
add_text(slide, .8, 5.25, 11.3, .85, "Current status: circular boundary work is in place; the interior logic still needs a dedicated horizontal-row pass.", 17, NAVY, True)
add_text(slide, .8, 6.22, 11.3, .42, "The saved artifact set has an explicitly named on-segment 150ft preview; the 24ft and 65ft panels use the nearest saved current previews.", 9, GREY)
add_text(slide, 12.35, 7.02, .35, .18, "3/6", 8, GREY)

# Slide 4: circular sweep update
slide = prs.slides[3]
clear(slide, 0)
heading(slide, "CIRCULAR SWEEP FIX", "Compact boundary coverage: two rows max", 4)
add_box(slide, .72, 1.85, 5.45, 3.95, RGBColor(255, 255, 255), LINE)
add_box(slide, 7.05, 1.85, 5.45, 3.95, RGBColor(255, 255, 255), LINE)
add_text(slide, 1.05, 2.15, 4.5, .34, "Circular sweep spacing", 18, BLUE, True)
add_text(slide, 1.05, 2.82, 4.7, 1.6, "Neighboring circular profiles are spaced closer together to create a small real overlap.\n\nThe planner now allows only 1–2 circular sweeps per tank; no rows beyond 2 are added.", 15)
add_text(slide, 7.38, 2.15, 4.5, .34, "Interior coverage status", 18, ORANGE, True)
add_text(slide, 7.38, 2.82, 4.7, 1.6, "Vertical plate coverage is the current interior focus.\n\nHorizontal plate rows are not implemented yet and remain the next geometry task.", 15)
add_text(slide, .75, 6.28, 11.2, .42, "Main takeaway: the boundary strategy is simpler and more controlled; the next gain comes from modeling horizontal weld rows.", 17, NAVY, True)
add_text(slide, 12.35, 7.02, .35, .18, "4/6", 8, GREY)

# Replace the former two-column status content with a focused 65ft circular-sweep preview.
clear(slide, 0)
heading(slide, "CIRCULAR SWEEP FIX", "Compact boundary coverage: two rows max", 4)
add_text(slide, .78, 1.78, 4.7, .34, "What changed", 18, BLUE, True)
add_text(slide, .78, 2.42, 4.72, 1.1, "The planner is capped at two circular sweep rows, so it does not keep adding tight inner rings.", 17, NAVY, True)
add_box(slide, .78, 3.82, 4.62, .01, LINE, LINE)
add_text(slide, .78, 4.12, 4.72, 1.1, "Neighboring rainbow footprints are packed closer using measured polygon overlap, leaving no visible gaps around the boundary.", 17, NAVY, True)
add_text(slide, .78, 5.86, 4.72, .35, "65ft example: two compact circular rows", 11, TEAL, True)
add_box(slide, 5.78, 1.7, 6.78, 4.98, RGBColor(255, 255, 255), LINE)
add_picture(slide, "after_65ft_path_planner.png", 5.83, 1.75, 6.68, 4.88)

# Slide 5: future work
slide = prs.slides[4]
clear(slide, 0)
heading(slide, "NEXT ENGINEERING RISKS", "Future work", 5)
add_box(slide, .72, 1.85, 5.45, 3.95, RGBColor(255, 255, 255), LINE)
add_box(slide, 7.05, 1.85, 5.45, 3.95, RGBColor(255, 255, 255), LINE)
add_text(slide, 1.05, 2.15, 4.8, .34, "Planning quality", 18, BLUE, True)
add_text(slide, 1.05, 2.82, 4.75, 1.8, "Implement horizontal-row detection and 90° profiles\n\nAdd Reeds–Shepp path planning for realistic robot motion\n\nKeep coverage and overlap checks geometry-based", 15)
add_text(slide, 7.38, 2.15, 4.8, .34, "Visualization and validation", 18, ORANGE, True)
add_text(slide, 7.38, 2.82, 4.75, 1.8, "Use the Reeds–Shepp path to connect scan poses\n\nVisualize the complete mission as an animation\n\nValidate motion, tether constraints, and field behavior", 15)
add_text(slide, 1.05, 6.28, 11, .42, "Main takeaway: the planner is moving from static coverage previews toward executable, visualized robot missions.", 17, NAVY, True)
add_text(slide, 12.35, 7.02, .35, .18, "5/6", 8, GREY)

# Slide 6: DXF completion outputs
slide = prs.slides.add_slide(prs.slide_layouts[6])
heading(slide, "DXF COMPLETER", "Partial observations -> completed layouts", 6)
add_text(
    slide, .75, 1.53, 11.35, .45,
    "One circular sweep preserves only observed weld fragments. The predictor normalizes that geometry by tank center and radius, selects the closest structural design family, and reconstructs the missing weld layout.",
    11.5, GREY,
)
completion = [
    ("24ft", BLUE),
    ("65ft", TEAL),
    ("150ft", ORANGE),
]
for index, (tank, color) in enumerate(completion):
    x = .55 + index * 4.18
    add_tag(slide, x + .12, 2.1, tank, color)
    add_text(slide, x + 1.47, 2.15, 2.0, .22, "observed", 10, GREY, True, PP_ALIGN.CENTER)
    add_box(slide, x, 2.42, 3.75, 1.66, RGBColor(255, 255, 255), LINE)
    add_picture(slide, f"dxf_prediction_validation/checkin/partial_{tank}.png", x + .78, 2.45, 2.18, 1.60)
    add_text(slide, x + 1.47, 4.19, 2.0, .22, "completed", 10, color, True, PP_ALIGN.CENTER)
    add_box(slide, x, 4.45, 3.75, 1.66, RGBColor(255, 255, 255), LINE)
    add_picture(slide, f"dxf_prediction_validation/checkin/completed_{tank}.png", x + .78, 4.48, 2.18, 1.60)
add_text(slide, .75, 6.45, 11.35, .30, "The comparison is based on geometry patterns, not a filename or literal tank size; the system can return ambiguous instead of completing when the observed evidence is insufficient.", 10.5, NAVY, True)

prs.save(OUT)
print(OUT)
